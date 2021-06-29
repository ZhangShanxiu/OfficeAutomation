#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Time    : 2021/6/19 16:37
# @Author  : Zhang Shanxiu
import os, cv2
from pptx import Presentation
from pptx.util import Cm, Pt


def main():
    opt = input('generate short?1/0: ')
    if opt == '0':
        path = './class'
    else:
        path = './image'
    file = Presentation('./class/result.pptx')
    # file.slides._sldIdLst.clear()
    # 记录页码
    page = 1
    for dirs in os.listdir(path=path):
        # A4纸张
        # file.slide_width = Cm(19.05)
        # file.slide_height = Cm(27.517)
        # print(dirs)
        file_path = os.path.join(path, dirs)
        print(file_path)
        if os.path.isdir(file_path):
            # cover_holder = 10
            # slide = file.slides.add_slide(file.slide_layouts[4])
            # with open(os.path.join(file_path, 'cover.txt'), encoding='UTF-8') as f:
            #     for line in f.readlines():
            #         slide.placeholders[cover_holder].text = line
            #         cover_holder += 1
            with open(os.path.join(file_path, 'readme.txt'), encoding='UTF-8') as f:
                lines = f.readlines()
            names = lines[0].split()
            file_number = len(names)

            if file_number <= 6:
                # 文件小于6个是一种版式
                template_number = 5
                picture_holder = 20
                yemei_holder = 26
                yejiao_holder = 27
            else:
                template_number = 6
                picture_holder = 22
                yemei_holder = 30
                yejiao_holder = 31
            slide = file.slides.add_slide(file.slide_layouts[template_number])
            # 页眉
            slide.placeholders[yemei_holder].text = 'RCAR2021 Technical Sessions'
            slide.placeholders[yejiao_holder].text = 'T-%d' % page
            page += 1
            place_holder = 11
            for i in range(1, len(lines)):
                slide.placeholders[place_holder].text = lines[i]
                place_holder += 1

            for name in names:
                picture_path = os.path.join(file_path, name + '.JPG')
                slide.placeholders[picture_holder].insert_picture(picture_path)
                picture_holder += 1
    file.slides._sldIdLst.remove(list(file.slides._sldIdLst)[0])
    file.slides._sldIdLst.remove(list(file.slides._sldIdLst)[0])
    file.slides._sldIdLst.remove(list(file.slides._sldIdLst)[0])
    if opt == '0':
        file.save('./class/program_normal_final.pptx')
        print('保存 ppt 文件：', './class/program_normal_final.pptx')
    else:
        file.save('./class/short_final.pptx')
        print('保存 ppt 文件：', './class/short_final.pptx')
    # file.ExportAsFixedFormat('./class/program.pdf', 2, PrintRange=None)
    # print('保存 pdf 文件：', './class/program.pdf')
                    

if __name__ == '__main__':
    main()